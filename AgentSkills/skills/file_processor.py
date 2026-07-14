"""
文件处理工具模块

提供以下功能：
- identify_file: 识别常见文件类型并提取基本信息
- modify_docx: 修改 Word (.docx) 文件（替换文本、添加段落/表格、修改标题）
- modify_pptx: 修改 PowerPoint (.pptx) 文件（替换文本、添加幻灯片、修改标题）
- modify_xlsx: 修改 Excel (.xlsx) 文件（修改单元格、追加行、重命名工作表、删除行）
"""

import os
import subprocess
import sys
import tempfile
from pathlib import Path


def _ensure_dependency(package_name: str, import_name: str = None) -> bool:
    """确保依赖包已安装，若未安装则自动通过 pip 安装

    Args:
        package_name: pip 包名
        import_name: import 时的模块名（默认为包名）

    Returns:
        是否安装成功
    """
    if import_name is None:
        import_name = package_name

    try:
        __import__(import_name)
        return True
    except ImportError:
        print(f"📦 正在自动安装 {package_name}...")
        try:
            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", package_name, "-q"],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            __import__(import_name)
            print(f"✅ {package_name} 安装成功")
            return True
        except Exception as e:
            print(f"❌ 自动安装 {package_name} 失败: {e}")
            return False


# ──────────────────────────────────────────────────────────
# A) 文件识别函数
# ──────────────────────────────────────────────────────────


def identify_file(file_path: str) -> str:
    """识别指定文件并提取基本信息。

    支持的格式：
    - 图片文件（jpg/jpeg/png/gif/bmp/webp/tiff）：提取宽高、格式、大小
    - Word 文档（.docx）：提取段落数、表格数、字符数
    - PowerPoint 演示文稿（.pptx）：提取幻灯片数、标题列表
    - Excel 工作簿（.xlsx）：提取工作表名、行列数
    - PDF 文件（.pdf）：提取页数、标题、作者
    - 文本文件（.txt / .csv / .md 等）：提取行数、字符数、编码

    Args:
        file_path: 待识别文件的路径

    Returns:
        格式化的文件信息字符串，如果出错则返回错误描述
    """
    path = Path(file_path)

    if not path.exists():
        return f"❌ 文件不存在: {file_path}"
    if not path.is_file():
        return f"❌ 路径不是文件: {file_path}"

    file_size = path.stat().st_size
    size_str = _format_size(file_size)
    ext = path.suffix.lower()

    try:
        # ── 图片文件 ──
        if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif"):
            return _identify_image(path, ext, size_str)

        # ── Word 文档 ──
        elif ext == ".docx":
            return _identify_docx(path, size_str)

        # ── PowerPoint 演示文稿 ──
        elif ext == ".pptx":
            return _identify_pptx(path, size_str)

        # ── Excel 工作簿 ──
        elif ext == ".xlsx":
            return _identify_xlsx(path, size_str)

        # ── PDF 文档 ──
        elif ext == ".pdf":
            return _identify_pdf(path, size_str)

        # ── 文本类文件 ──
        elif ext in (".txt", ".csv", ".md", ".json", ".xml", ".yaml", ".yml",
                     ".log", ".ini", ".cfg", ".conf", ".toml", ".env",
                     ".py", ".js", ".ts", ".go", ".c", ".cpp", ".h", ".hpp",
                     ".java", ".rs", ".rb", ".php", ".html", ".css", ".sh",
                     ".sql", ".r", ".swift", ".kt", ".dart"):
            return _identify_text_file(path, ext, size_str)

        # ── 其他文件 ──
        else:
            return (
                f"📄 文件: {path.name}\n"
                f"├─ 类型: 未知格式 ({ext or '无后缀'})\n"
                f"├─ 大小: {size_str}\n"
                f"├─ 路径: {path.resolve()}\n"
                f"└─ 提示: 该格式未实现专门的识别逻辑"
            )

    except Exception as e:
        return f"❌ 识别文件时出错: {e}"


def _format_size(size_bytes: int) -> str:
    """将字节数转换为可读的大小字符串"""
    for unit in ("B", "KB", "MB", "GB"):
        if size_bytes < 1024:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} TB"


def _identify_image(path: Path, ext: str, size_str: str) -> str:
    """识别图片文件"""
    if not _ensure_dependency("Pillow", "PIL"):
        return (
            f"🖼️ 文件: {path.name}\n"
            f"├─ 类型: 图片 ({ext})\n"
            f"├─ 大小: {size_str}\n"
            f"└─ ⚠️ 无法加载 Pillow 库，无法获取详细图片信息"
        )
    from PIL import Image

    img = Image.open(path)
    width, height = img.size
    mode = img.mode
    fmt = img.format or ext.upper().lstrip(".")
    info_str = (
        f"🖼️ 文件: {path.name}\n"
        f"├─ 类型: 图片\n"
        f"├─ 格式: {fmt}\n"
        f"├─ 尺寸: {width} × {height} 像素\n"
        f"├─ 色彩模式: {mode}\n"
        f"├─ 大小: {size_str}\n"
        f"└─ 路径: {path.resolve()}"
    )
    return info_str


def _identify_docx(path: Path, size_str: str) -> str:
    """识别 Word 文档"""
    if not _ensure_dependency("python-docx", "docx"):
        return (
            f"📄 文件: {path.name}\n"
            f"├─ 类型: Word 文档\n"
            f"├─ 大小: {size_str}\n"
            f"└─ ⚠️ 无法加载 python-docx 库，无法获取详细文档信息"
        )
    from docx import Document

    doc = Document(str(path))
    para_count = len(doc.paragraphs)
    table_count = len(doc.tables)
    # 统计非空段落数
    non_empty_paras = sum(1 for p in doc.paragraphs if p.text.strip())
    # 统计总字符数
    total_chars = sum(len(p.text) for p in doc.paragraphs)

    # 获取文档标题（第一个标题样式段落）
    title_text = ""
    for p in doc.paragraphs:
        if p.style and "Title" in p.style.name:
            title_text = p.text.strip()
            break

    info_str = (
        f"📄 文件: {path.name}\n"
        f"├─ 类型: Word 文档 (.docx)\n"
        f"├─ 段落数: {para_count}（非空: {non_empty_paras}）\n"
        f"├─ 表格数: {table_count}\n"
        f"├─ 字符数: {total_chars}\n"
        f"├─ 大小: {size_str}\n"
        f"├─ 路径: {path.resolve()}"
    )
    if title_text:
        info_str += f"\n└─ 文档标题: {title_text}"
    return info_str


def _identify_pptx(path: Path, size_str: str) -> str:
    """识别 PowerPoint 演示文稿"""
    if not _ensure_dependency("python-pptx"):
        return (
            f"📑 文件: {path.name}\n"
            f"├─ 类型: PowerPoint 演示文稿\n"
            f"├─ 大小: {size_str}\n"
            f"└─ ⚠️ 无法加载 python-pptx 库，无法获取详细幻灯片信息"
        )
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_count = len(prs.slides)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    slide_titles = []
    for i, slide in enumerate(prs.slides, start=1):
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
        slide_titles.append(f"    {i}. {title_text or '(无标题)'}")

    title_list_str = "\n".join(slide_titles[:10])  # 最多显示10张
    if len(slide_titles) > 10:
        title_list_str += f"\n    ... 还有 {len(slide_titles) - 10} 张幻灯片"

    info_str = (
        f"📑 文件: {path.name}\n"
        f"├─ 类型: PowerPoint 演示文稿 (.pptx)\n"
        f"├─ 幻灯片数: {slide_count}\n"
        f"├─ 幻灯片尺寸: {slide_width} × {slide_height} (EMU)\n"
        f"├─ 大小: {size_str}\n"
        f"├─ 路径: {path.resolve()}\n"
        f"└─ 幻灯片标题:\n{title_list_str}"
    )
    return info_str


def _identify_xlsx(path: Path, size_str: str) -> str:
    """识别 Excel 工作簿"""
    if not _ensure_dependency("openpyxl"):
        return (
            f"📊 文件: {path.name}\n"
            f"├─ 类型: Excel 工作簿\n"
            f"├─ 大小: {size_str}\n"
            f"└─ ⚠️ 无法加载 openpyxl 库，无法获取详细工作表信息"
        )
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheet_count = len(wb.sheetnames)

    sheets_info = []
    for name in wb.sheetnames:
        ws = wb[name]
        max_row = ws.max_row or 0
        max_col = ws.max_column or 0
        sheets_info.append(f"    📊 {name} ({max_row} 行 × {max_col} 列)")

    sheet_list_str = "\n".join(sheets_info)
    wb.close()

    info_str = (
        f"📊 文件: {path.name}\n"
        f"├─ 类型: Excel 工作簿 (.xlsx)\n"
        f"├─ 工作表数: {sheet_count}\n"
        f"├─ 大小: {size_str}\n"
        f"├─ 路径: {path.resolve()}\n"
        f"└─ 工作表详情:\n{sheet_list_str}"
    )
    return info_str


def _identify_pdf(path: Path, size_str: str) -> str:
    """识别 PDF 文档"""
    info = (
        f"📕 文件: {path.name}\n"
        f"├─ 类型: PDF 文档\n"
        f"├─ 大小: {size_str}\n"
        f"├─ 路径: {path.resolve()}\n"
    )
    # 尝试获取页数（可通过 pdfminer/pdfplumber 或命令行工具）
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            page_count = len(pdf.pages)
            info += f"└─ 页数: {page_count}"
    except ImportError:
        try:
            import PyPDF2
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                page_count = len(reader.pages)
                info += f"└─ 页数: {page_count}"
        except ImportError:
            info += "└─ 页数: (需要安装 pdfplumber 或 PyPDF2 以获取)"
    except Exception:
        info += "└─ 页数: (无法读取)"
    return info


def _identify_text_file(path: Path, ext: str, size_str: str) -> str:
    """识别文本文件"""
    ext_map = {
        ".txt": "纯文本", ".csv": "CSV 表格", ".md": "Markdown 文档",
        ".json": "JSON 数据", ".xml": "XML 文档", ".yaml": "YAML 配置",
        ".yml": "YAML 配置", ".toml": "TOML 配置", ".env": "环境变量文件",
        ".ini": "INI 配置文件", ".cfg": "配置文件", ".conf": "配置文件",
        ".log": "日志文件",
        ".py": "Python 源码", ".js": "JavaScript 源码", ".ts": "TypeScript 源码",
        ".go": "Go 源码", ".c": "C 源码", ".cpp": "C++ 源码",
        ".h": "C/C++ 头文件", ".hpp": "C++ 头文件",
        ".java": "Java 源码", ".rs": "Rust 源码", ".rb": "Ruby 源码",
        ".php": "PHP 源码", ".html": "HTML 文档", ".css": "CSS 样式表",
        ".sh": "Shell 脚本", ".sql": "SQL 脚本", ".r": "R 脚本",
        ".swift": "Swift 源码", ".kt": "Kotlin 源码", ".dart": "Dart 源码",
    }
    file_type = ext_map.get(ext, f"文本文件 ({ext})")

    # 尝试检测编码并读取
    encoding = "utf-8"
    encodings_to_try = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1", "shift-jis", "big5"]
    content = None
    for enc in encodings_to_try:
        try:
            with open(path, "r", encoding=enc) as f:
                content = f.read()
            encoding = enc
            break
        except (UnicodeDecodeError, UnicodeError):
            continue

    if content is None:
        try:
            with open(path, "rb") as f:
                raw = f.read(100)
            encoding = f"binary (无法解码, 前 {len(raw)} 字节: {raw[:40].hex()})"
            lines = 0
            chars = 0
        except Exception:
            encoding = "unknown"
            lines = 0
            chars = 0
    else:
        lines = content.count("\n") + (1 if content and not content.endswith("\n") else 0)
        if not content:
            lines = 0
        chars = len(content)

    info_str = (
        f"📄 文件: {path.name}\n"
        f"├─ 类型: {file_type}\n"
        f"├─ 行数: {lines}\n"
        f"├─ 字符数: {chars}\n"
        f"├─ 编码: {encoding}\n"
        f"├─ 大小: {size_str}\n"
        f"└─ 路径: {path.resolve()}"
    )
    return info_str


# ──────────────────────────────────────────────────────────
# B) DOCX 修改函数
# ──────────────────────────────────────────────────────────


def modify_docx(file_path: str, operations: list, output_path: str = None) -> str:
    """修改 Word (.docx) 文件。

    支持的操作类型：
        - {"type": "replace_text", "old": "原文本", "new": "新文本"}
            替换文档中所有匹配的文本
        - {"type": "add_paragraph", "text": "段落内容", "style": "Normal"}
            在文档末尾添加段落（style 可选，默认 "Normal"）
        - {"type": "add_table", "rows": 3, "cols": 4, "data": [["a","b","c","d"],...]}
            在文档末尾添加表格（data 可选，不提供则创建空表）
        - {"type": "set_title", "text": "新标题"}
            修改文档标题（找到样式为 Title 的段落并修改其文本）

    Args:
        file_path: 源 .docx 文件路径
        operations: 操作列表，每个元素为一个 dict
        output_path: 输出文件路径（默认覆盖原文件）

    Returns:
        操作结果描述字符串
    """
    if not _ensure_dependency("python-docx", "docx"):
        return "❌ 无法加载 python-docx 库，请手动安装: pip install python-docx"

    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    if not os.path.isfile(file_path):
        return f"❌ 文件不存在: {file_path}"

    output = output_path or file_path

    try:
        doc = Document(file_path)
        results = []

        for i, op in enumerate(operations, start=1):
            op_type = op.get("type", "")

            if op_type == "replace_text":
                old_text = op.get("old", "")
                new_text = op.get("new", "")
                if not old_text:
                    results.append(f"  [{i}] ⚠️ replace_text: 缺少 'old' 参数，跳过")
                    continue
                count = 0
                for paragraph in doc.paragraphs:
                    if old_text in paragraph.text:
                        for run in paragraph.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
                                count += 1
                # 也替换表格中的文本
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                if old_text in paragraph.text:
                                    for run in paragraph.runs:
                                        if old_text in run.text:
                                            run.text = run.text.replace(old_text, new_text)
                                            count += 1
                results.append(f"  [{i}] ✅ replace_text: 替换了 {count} 处 \"{old_text}\" → \"{new_text}\"")

            elif op_type == "add_paragraph":
                text = op.get("text", "")
                style_name = op.get("style", "Normal")
                p = doc.add_paragraph(text)
                try:
                    p.style = doc.styles[style_name]
                except KeyError:
                    results.append(f"  [{i}] ⚠️ add_paragraph: 样式 \"{style_name}\" 不存在，使用默认样式")
                results.append(f"  [{i}] ✅ add_paragraph: 添加了段落 \"{text[:50]}{'...' if len(text) > 50 else ''}\"")

            elif op_type == "add_table":
                rows = op.get("rows", 1)
                cols = op.get("cols", 1)
                data = op.get("data", None)
                if data:
                    table = doc.add_table(rows=len(data), cols=len(data[0]) if data else 1)
                    for r_idx, row_data in enumerate(data):
                        for c_idx, cell_value in enumerate(row_data):
                            table.cell(r_idx, c_idx).text = str(cell_value)
                else:
                    table = doc.add_table(rows=rows, cols=cols)
                # 应用表格样式
                try:
                    table.style = "Table Grid"
                except KeyError:
                    pass
                results.append(f"  [{i}] ✅ add_table: 添加了 {rows}×{cols} 表格" +
                               (f" (填充了 {len(data)} 行数据)" if data else ""))

            elif op_type == "set_title":
                new_title = op.get("text", "")
                if not new_title:
                    results.append(f"  [{i}] ⚠️ set_title: 缺少 'text' 参数，跳过")
                    continue
                found = False
                for p in doc.paragraphs:
                    if p.style and "Title" in p.style.name:
                        # 清除原有 runs 并设置新文本
                        p.clear()
                        run = p.add_run(new_title)
                        found = True
                        break
                if not found:
                    # 如果文档没有 Title 样式的段落，在开头插入一个
                    p = doc.paragraphs[0].insert_paragraph_before(new_title)
                    try:
                        p.style = doc.styles["Title"]
                    except KeyError:
                        pass
                    results.append(f"  [{i}] ℹ️ set_title: 文档无标题段落，已在开头添加新标题")
                results.append(f"  [{i}] ✅ set_title: 文档标题已设置为 \"{new_title}\"")

            else:
                results.append(f"  [{i}] ⚠️ 未知操作类型: \"{op_type}\"，跳过")

        doc.save(output)
        summary = "\n".join(results)
        path_info = f"（覆盖原文件）" if output == file_path else f"→ {output}"
        return f"✅ docx 修改完成 {path_info}\n{summary}"

    except Exception as e:
        return f"❌ 修改 docx 文件时出错: {e}"


# ──────────────────────────────────────────────────────────
# C) PPTX 修改函数
# ──────────────────────────────────────────────────────────


def modify_pptx(file_path: str, operations: list, output_path: str = None) -> str:
    """修改 PowerPoint (.pptx) 文件。

    支持的操作类型：
        - {"type": "replace_text", "old": "原文本", "new": "新文本", "slide_index": 0}
            替换指定幻灯片中的所有文本（slide_index 从 0 开始，不指定则替换所有幻灯片）
        - {"type": "add_slide", "title": "幻灯片标题", "content": ["项目1","项目2"]}
            在末尾添加新幻灯片（content 为正文文本列表）
        - {"type": "set_slide_title", "slide_index": 0, "text": "新标题"}
            修改指定幻灯片的标题文本

    Args:
        file_path: 源 .pptx 文件路径
        operations: 操作列表，每个元素为一个 dict
        output_path: 输出文件路径（默认覆盖原文件）

    Returns:
        操作结果描述字符串
    """
    if not _ensure_dependency("python-pptx"):
        return "❌ 无法加载 python-pptx 库，请手动安装: pip install python-pptx"

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    if not os.path.isfile(file_path):
        return f"❌ 文件不存在: {file_path}"

    output = output_path or file_path

    try:
        prs = Presentation(file_path)
        results = []

        for i, op in enumerate(operations, start=1):
            op_type = op.get("type", "")

            if op_type == "replace_text":
                old_text = op.get("old", "")
                new_text = op.get("new", "")
                slide_index = op.get("slide_index", None)

                if not old_text:
                    results.append(f"  [{i}] ⚠️ replace_text: 缺少 'old' 参数，跳过")
                    continue

                count = 0
                slides_to_process = []
                if slide_index is not None:
                    if 0 <= slide_index < len(prs.slides):
                        slides_to_process = [prs.slides[slide_index]]
                    else:
                        results.append(f"  [{i}] ⚠️ replace_text: slide_index {slide_index} 超出范围（共 {len(prs.slides)} 张幻灯片）")
                        continue
                else:
                    slides_to_process = list(prs.slides)

                for slide in slides_to_process:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text)
                                        count += 1

                location = f"幻灯片 {slide_index}" if slide_index is not None else "所有幻灯片"
                results.append(f"  [{i}] ✅ replace_text: 在 {location} 中替换了 {count} 处 \"{old_text}\" → \"{new_text}\"")

            elif op_type == "add_slide":
                title_text = op.get("title", "新幻灯片")
                content_items = op.get("content", [])

                # 添加空白幻灯片（使用第一个版式）
                slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None
                if slide_layout is None:
                    results.append(f"  [{i}] ❌ add_slide: 无法获取幻灯片版式")
                    continue

                slide = prs.slides.add_slide(slide_layout)

                # 设置标题
                if slide.shapes.title:
                    slide.shapes.title.text = title_text

                # 添加内容（如果版式有正文占位符，使用它；否则添加文本框）
                if content_items:
                    # 查找正文占位符
                    body_shape = None
                    for shape in slide.shapes:
                        if shape != slide.shapes.title and shape.has_text_frame:
                            body_shape = shape
                            break

                    if body_shape:
                        tf = body_shape.text_frame
                        tf.text = ""
                        for idx, item in enumerate(content_items):
                            if idx == 0:
                                tf.text = str(item)
                            else:
                                p = tf.add_paragraph()
                                p.text = str(item)
                    else:
                        # 在右侧添加文本框
                        left = Inches(1)
                        top = Inches(2)
                        width = Inches(8)
                        height = Inches(5)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = ""
                        for idx, item in enumerate(content_items):
                            if idx == 0:
                                tf.text = str(item)
                            else:
                                p = tf.add_paragraph()
                                p.text = str(item)

                content_summary = f" ({len(content_items)} 个要点)" if content_items else ""
                results.append(f"  [{i}] ✅ add_slide: 添加了幻灯片 \"{title_text}\"{content_summary}")

            elif op_type == "set_slide_title":
                slide_index = op.get("slide_index", -1)
                new_title = op.get("text", "")

                if not new_title:
                    results.append(f"  [{i}] ⚠️ set_slide_title: 缺少 'text' 参数，跳过")
                    continue

                if slide_index < 0 or slide_index >= len(prs.slides):
                    results.append(f"  [{i}] ⚠️ set_slide_title: slide_index {slide_index} 超出范围（共 {len(prs.slides)} 张幻灯片）")
                    continue

                slide = prs.slides[slide_index]
                if slide.shapes.title:
                    slide.shapes.title.text = new_title
                    results.append(f"  [{i}] ✅ set_slide_title: 幻灯片 {slide_index} 标题已设置为 \"{new_title}\"")
                else:
                    results.append(f"  [{i}] ⚠️ set_slide_title: 幻灯片 {slide_index} 没有标题占位符")

            else:
                results.append(f"  [{i}] ⚠️ 未知操作类型: \"{op_type}\"，跳过")

        prs.save(output)
        summary = "\n".join(results)
        path_info = f"（覆盖原文件）" if output == file_path else f"→ {output}"
        return f"✅ pptx 修改完成 {path_info}\n{summary}"

    except Exception as e:
        return f"❌ 修改 pptx 文件时出错: {e}"


# ──────────────────────────────────────────────────────────
# D) XLSX 修改函数
# ──────────────────────────────────────────────────────────


def modify_xlsx(file_path: str, operations: list, output_path: str = None) -> str:
    """修改 Excel (.xlsx) 文件。

    支持的操作类型：
        - {"type": "set_cell", "sheet": "Sheet1", "row": 1, "col": 1, "value": "新值"}
            修改指定单元格的值（row 和 col 从 1 开始计数）
        - {"type": "add_row", "sheet": "Sheet1", "data": ["a","b","c"]}
            在指定工作表的末尾追加一行
        - {"type": "set_sheet_name", "sheet": "Sheet1", "new_name": "新名称"}
            重命名工作表
        - {"type": "delete_row", "sheet": "Sheet1", "row_index": 2}
            删除指定行（row_index 从 1 开始计数，删除表头行请谨慎）

    Args:
        file_path: 源 .xlsx 文件路径
        operations: 操作列表，每个元素为一个 dict
        output_path: 输出文件路径（默认覆盖原文件）

    Returns:
        操作结果描述字符串
    """
    if not _ensure_dependency("openpyxl"):
        return "❌ 无法加载 openpyxl 库，请手动安装: pip install openpyxl"

    import openpyxl

    if not os.path.isfile(file_path):
        return f"❌ 文件不存在: {file_path}"

    output = output_path or file_path

    try:
        wb = openpyxl.load_workbook(file_path)
        results = []

        for i, op in enumerate(operations, start=1):
            op_type = op.get("type", "")

            if op_type == "set_cell":
                sheet_name = op.get("sheet", "Sheet1")
                row = op.get("row")
                col = op.get("col")
                value = op.get("value", "")

                if row is None or col is None:
                    results.append(f"  [{i}] ⚠️ set_cell: 缺少 'row' 或 'col' 参数，跳过")
                    continue

                if sheet_name not in wb.sheetnames:
                    results.append(f"  [{i}] ⚠️ set_cell: 工作表 \"{sheet_name}\" 不存在，跳过")
                    continue

                ws = wb[sheet_name]
                ws.cell(row=row, column=col, value=value)
                results.append(f"  [{i}] ✅ set_cell: [{sheet_name}] ({row},{col}) = \"{value}\"")

            elif op_type == "add_row":
                sheet_name = op.get("sheet", "Sheet1")
                data = op.get("data", [])

                if sheet_name not in wb.sheetnames:
                    results.append(f"  [{i}] ⚠️ add_row: 工作表 \"{sheet_name}\" 不存在，跳过")
                    continue

                ws = wb[sheet_name]
                ws.append(data)
                results.append(f"  [{i}] ✅ add_row: 在 [{sheet_name}] 末尾添加了 {len(data)} 列数据")

            elif op_type == "set_sheet_name":
                sheet_name = op.get("sheet", "")
                new_name = op.get("new_name", "")

                if not sheet_name:
                    results.append(f"  [{i}] ⚠️ set_sheet_name: 缺少 'sheet' 参数，跳过")
                    continue
                if not new_name:
                    results.append(f"  [{i}] ⚠️ set_sheet_name: 缺少 'new_name' 参数，跳过")
                    continue
                if sheet_name not in wb.sheetnames:
                    results.append(f"  [{i}] ⚠️ set_sheet_name: 工作表 \"{sheet_name}\" 不存在，跳过")
                    continue
                if new_name in wb.sheetnames:
                    results.append(f"  [{i}] ⚠️ set_sheet_name: 目标名称 \"{new_name}\" 已存在，跳过")
                    continue

                ws = wb[sheet_name]
                ws.title = new_name
                results.append(f"  [{i}] ✅ set_sheet_name: \"{sheet_name}\" → \"{new_name}\"")

            elif op_type == "delete_row":
                sheet_name = op.get("sheet", "Sheet1")
                row_index = op.get("row_index")

                if row_index is None:
                    results.append(f"  [{i}] ⚠️ delete_row: 缺少 'row_index' 参数，跳过")
                    continue

                if sheet_name not in wb.sheetnames:
                    results.append(f"  [{i}] ⚠️ delete_row: 工作表 \"{sheet_name}\" 不存在，跳过")
                    continue

                ws = wb[sheet_name]
                if row_index < 1 or row_index > (ws.max_row or 0):
                    results.append(f"  [{i}] ⚠️ delete_row: row_index {row_index} 超出范围（共 {ws.max_row} 行），跳过")
                    continue

                ws.delete_rows(row_index, 1)
                results.append(f"  [{i}] ✅ delete_row: 从 [{sheet_name}] 中删除了第 {row_index} 行")

            else:
                results.append(f"  [{i}] ⚠️ 未知操作类型: \"{op_type}\"，跳过")

        wb.save(output)
        summary = "\n".join(results)
        path_info = f"（覆盖原文件）" if output == file_path else f"→ {output}"
        return f"✅ xlsx 修改完成 {path_info}\n{summary}"

    except Exception as e:
        return f"❌ 修改 xlsx 文件时出错: {e}"
